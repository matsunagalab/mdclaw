"""Generate MDClaw progress-report slides from scratch (no template)."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

OUTPUT = Path(
    "/Users/yasu/gdrive/work/mdclaw/slide/"
    "20260420_進捗報告_MDClawによるエージェント駆動MD_松永.pptx"
)

# ---------- palette ----------
INK = RGBColor(0x0F, 0x17, 0x2A)         # near-black
INK_MUTED = RGBColor(0x4B, 0x55, 0x63)   # slate-600
INK_DIM = RGBColor(0x9A, 0xA2, 0xAE)     # slate-400
LINE = RGBColor(0xE2, 0xE8, 0xF0)        # slate-200
BG = RGBColor(0xFF, 0xFF, 0xFF)          # white
PANEL = RGBColor(0xF8, 0xFA, 0xFC)       # slate-50
ACCENT = RGBColor(0x0F, 0x76, 0x6E)      # teal-700 — main accent
ACCENT_DARK = RGBColor(0x0B, 0x4F, 0x4A)
ACCENT_LIGHT = RGBColor(0xD6, 0xEE, 0xEA)

# node colors, semantic
C_FETCH = RGBColor(0x38, 0xBD, 0xF8)     # sky-400
C_STEP = RGBColor(0x2B, 0x73, 0xE6)      # blue-600
C_EQ = RGBColor(0x10, 0xB9, 0x81)        # emerald-500
C_PROD = RGBColor(0x8B, 0x5C, 0xF6)      # violet-500
C_AGENT = RGBColor(0xDC, 0x26, 0x26)     # red-600
C_EVENT = RGBColor(0x6B, 0x72, 0x80)     # gray-500
C_WARN = RGBColor(0xF5, 0x9E, 0x0B)      # amber-500
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# slide size
SLIDE_W = 13.333
SLIDE_H = 7.5

# layout constants
MARGIN_X = 0.55
TOP_BAR_H = 0.35      # colored context bar
TITLE_Y = 0.70
TITLE_H = 0.75
BODY_Y = 1.65
BODY_H = 5.30
FOOTER_Y = 7.10


# ---------- low-level helpers ----------

def _add_textbox(slide, x, y, w, h, fill=None, border=None, border_width_pt=0.75):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    return tb


def text(slide, x, y, w, h, content, size=14, bold=False, italic=False,
         color=INK, align=PP_ALIGN.LEFT, font="Helvetica Neue",
         anchor=MSO_ANCHOR.TOP, line_spacing=1.15, space_after=0):
    """Single-run single-paragraph textbox."""
    tb = _add_textbox(slide, x, y, w, h)
    tf = tb.text_frame
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = content
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size=14, color=INK,
            font="Helvetica Neue", line_spacing=1.25, space_after=6,
            indent_char="•", indent_color=None):
    """Simple bullet list. Each item is a string (or tuple (text, indent_level))."""
    tb = _add_textbox(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            line, level = item
        else:
            line, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        if not line:
            # blank spacer paragraph — no bullet
            r = p.add_run()
            r.text = " "
            r.font.size = Pt(max(6, size // 2))
            continue
        # bullet
        marker = indent_char if level == 0 else "–"
        pad = "    " * level
        r_marker = p.add_run()
        r_marker.text = f"{pad}{marker}  "
        r_marker.font.size = Pt(size)
        r_marker.font.name = font
        r_marker.font.color.rgb = indent_color or ACCENT
        r_marker.font.bold = True
        # text
        r_text = p.add_run()
        r_text.text = line
        r_text.font.size = Pt(size)
        r_text.font.name = font
        r_text.font.color.rgb = color
    return tb


def filled_rect(slide, x, y, w, h, fill, line=None, line_width_pt=0.0,
                dash=False, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None and line_width_pt == 0:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line if line is not None else INK
        sh.line.width = Pt(line_width_pt)
        if dash:
            sh.line.dash_style = 7
    # remove default shadow
    sp = sh.shadow
    try:
        sp.inherit = False
    except Exception:
        pass
    return sh


def node(slide, x, y, w, h, label, fill=C_STEP, text_color=WHITE,
         size=11, bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         font="Helvetica Neue"):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = fill
    sh.line.width = Pt(0)
    tf = sh.text_frame
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color
    r.font.name = font
    return sh


def arrow(slide, x1, y1, x2, y2, color=INK_MUTED, width_pt=1.3,
          head="triangle"):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width_pt)
    ln = c.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", head)
    tail.set("w", "med")
    tail.set("len", "med")
    return c


def chrome(slide, context_label, slide_no, total):
    """Top color bar + page number + bottom hairline."""
    # colored top strip
    filled_rect(slide, 0, 0, SLIDE_W, TOP_BAR_H, ACCENT)
    # context label on left
    text(slide, MARGIN_X, 0.03, 5.0, TOP_BAR_H - 0.05,
         context_label.upper(), size=11, bold=True, color=WHITE,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # page number on right
    text(slide, SLIDE_W - MARGIN_X - 3.0, 0.03, 3.0, TOP_BAR_H - 0.05,
         f"{slide_no} / {total}   ·   MDClaw progress · 2026-04",
         size=10, color=ACCENT_LIGHT, align=PP_ALIGN.RIGHT,
         anchor=MSO_ANCHOR.MIDDLE)
    # bottom hairline
    filled_rect(slide, 0, 7.47, SLIDE_W, 0.03, LINE)


def slide_title(slide, title, subtitle=None):
    text(slide, MARGIN_X, TITLE_Y, SLIDE_W - 2 * MARGIN_X, TITLE_H,
         title, size=26, bold=True, color=INK,
         font="Helvetica Neue", line_spacing=1.1)
    if subtitle:
        text(slide, MARGIN_X, TITLE_Y + 0.85, SLIDE_W - 2 * MARGIN_X, 0.4,
             subtitle, size=13, italic=True, color=INK_MUTED,
             font="Helvetica Neue")


def add_slide(prs, context, title, subtitle=None):
    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)
    # remember meta — page number filled after all slides are added
    slide._mdclaw_context = context
    slide._mdclaw_title = title
    slide._mdclaw_subtitle = subtitle
    return slide


# ---------- figures ----------

def figure_dag_primer(slide, y0):
    """Linear chain with branching at both prep and eq — a real DAG."""
    nw, nh = 1.05, 0.55
    # main chain positions
    xs = [0.75, 1.95, 3.15, 4.35, 5.55, 6.75]
    main_nodes = [("source", C_FETCH), ("prep", C_STEP), ("solv", C_STEP),
                  ("topo", C_STEP), ("eq", C_EQ), ("prod_001", C_PROD)]
    for (lbl, col), x in zip(main_nodes, xs):
        node(slide, x, y0, nw, nh, lbl, fill=col)
    for i in range(len(xs) - 1):
        arrow(slide, xs[i] + nw, y0 + nh / 2,
              xs[i + 1], y0 + nh / 2, width_pt=1.5)
    # branch at prep: a mutant sibling chain going DOWN
    prep_x = xs[1]
    mut_y = y0 + 1.35
    mut_chain = [("prep_mut", C_STEP), ("solv_mut", C_STEP),
                 ("topo_mut", C_STEP), ("eq_mut", C_EQ),
                 ("prod_mut", C_PROD)]
    mut_xs = xs[1:]  # align under solv..prod_001
    for (lbl, col), x in zip(mut_chain, mut_xs):
        node(slide, x, mut_y, nw, nh, lbl, fill=col, size=10)
    # arrow from prep into prep_mut
    arrow(slide, prep_x + nw / 2, y0 + nh,
          mut_xs[0] + nw / 2, mut_y, width_pt=1.5, color=INK_MUTED)
    for i in range(len(mut_xs) - 1):
        arrow(slide, mut_xs[i] + nw, mut_y + nh / 2,
              mut_xs[i + 1], mut_y + nh / 2, width_pt=1.3)
    # branches from eq (xs[4]): prod_002, prod_003 above
    eq_x = xs[4]
    bx = 8.25
    node(slide, bx, y0 - 1.20, 1.65, nh, "prod_002  seed=42",
         fill=C_PROD, size=10)
    node(slide, bx, y0 - 0.50, 1.65, nh, "prod_003  T=320 K",
         fill=C_PROD, size=10)
    arrow(slide, eq_x + nw, y0 + nh / 2, bx, y0 - 1.20 + nh / 2,
          width_pt=1.5)
    arrow(slide, eq_x + nw, y0 + nh / 2, bx, y0 - 0.50 + nh / 2,
          width_pt=1.5)
    # annotations
    text(slide, 10.05, y0 - 1.28, 3.1, 0.3,
         "branch at eq → more prod", size=10, italic=True, color=INK_MUTED)
    text(slide, 10.05, y0 - 0.58, 3.1, 0.3,
         "(different seed / temperature)", size=10, italic=True, color=INK_DIM)
    text(slide, 8.10, mut_y + 0.10, 4.9, 0.3,
         "branch at prep → mutant system (shares topology code path)",
         size=10, italic=True, color=INK_MUTED)
    # callout for the word "branching"
    filled_rect(slide, 0.75, mut_y + 0.75, 6.2, 0.55, ACCENT_LIGHT,
                line=ACCENT, line_width_pt=0.75)
    text(slide, 0.9, mut_y + 0.83, 6.0, 0.4,
         "A DAG means branching is a first-class move, not a special case.",
         size=12, bold=True, color=ACCENT_DARK)


def figure_positioning(slide, y0):
    """Two-column comparison: LangGraph vs mdclaw."""
    label_w = 2.35
    col_w = (SLIDE_W - 2 * MARGIN_X - label_w - 0.3) / 2
    x_label = MARGIN_X
    xs = [MARGIN_X + label_w + 0.15,
          MARGIN_X + label_w + 0.15 + col_w + 0.15]
    # headers
    filled_rect(slide, xs[0], y0, col_w, 0.55, INK_DIM)
    filled_rect(slide, xs[1], y0, col_w, 0.55, ACCENT)
    text(slide, xs[0], y0 + 0.1, col_w, 0.4,
         "LangGraph-style control flow", size=13, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, xs[1], y0 + 0.1, col_w, 0.4,
         "mdclaw artifact lineage", size=13, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rows = [
        ("Graph represents", "Reasoning / control flow",
         "Lineage of computation artifacts"),
        ("Where state lives", "Library memory",
         "Filesystem (node.json, artifacts/, events/)"),
        ("Adding a node", "Code change to graph definition",
         "Single create_node call by the agent"),
        ("Agent location", "Inside the graph",
         "Outside the graph — operates on it"),
        ("Durable execution", "Hard",
         "Native — any process re-enters from disk"),
    ]
    row_h = 0.62
    for i, (label, left, right) in enumerate(rows):
        ry = y0 + 0.6 + i * row_h
        band = PANEL if i % 2 == 0 else WHITE
        filled_rect(slide, MARGIN_X, ry, SLIDE_W - 2 * MARGIN_X, row_h, band)
        text(slide, x_label + 0.15, ry + 0.1, label_w - 0.2, row_h - 0.2,
             label, size=12, bold=True, color=INK_MUTED,
             anchor=MSO_ANCHOR.MIDDLE)
        text(slide, xs[0] + 0.15, ry + 0.12, col_w - 0.3, row_h - 0.2,
             left, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, xs[1] + 0.15, ry + 0.12, col_w - 0.3, row_h - 0.2,
             right, size=12, color=INK, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)


def figure_dirtree(slide, x, y):
    """Directory tree of schema v3."""
    lines = [
        ("job_dir/",            0, True,  INK),
        ("progress.json",       1, False, INK_MUTED),
        ("# thin index of nodes + cached summary", 2, False, INK_DIM),
        ("nodes/",              1, True,  INK),
        ("source_001/",          2, True,  C_STEP),
        ("node.json · node.lock · artifacts/…", 3, False, INK_MUTED),
        ("prep_001/  solv_001/  topo_001/", 2, True, C_STEP),
        ("eq_001/",              2, True,  C_EQ),
        ("equilibrated.xml (saveState)", 3, False, INK_MUTED),
        ("prod_001/  prod_002/  prod_003/  branches", 2, True, C_PROD),
        ("trajectory.dcd · state.xml · checkpoint.chk", 3, False, INK_MUTED),
        ("events/",              1, True,  INK),
        ("<ISO8601>_<node>_<event>.json  (append-only)", 2, False, INK_MUTED),
    ]
    tb = _add_textbox(slide, x, y, 6.3, 4.6)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (line, level, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.1
        p.space_after = Pt(1)
        r = p.add_run()
        r.text = ("  " * level) + line
        r.font.size = Pt(12)
        r.font.name = "Menlo"
        r.font.bold = bold
        r.font.color.rgb = col


def figure_cli_groups(slide, y0):
    """Grouped server boxes."""
    groups = [
        ("research",  "download_structure · get_alphafold_structure · inspect_molecules · register_local_structure", C_STEP),
        ("structure", "prepare_complex · clean_protein · clean_ligand · run_antechamber_robust · create_mutated_structure", C_STEP),
        ("genesis",   "boltz2_protein_from_seq · rdkit_validate_smiles · pubchem_get_smiles_from_name", C_FETCH),
        ("solvation", "solvate_structure · embed_in_membrane · list_available_lipids", C_STEP),
        ("amber",     "build_amber_system  (openmmforcefields + Pablo)", C_STEP),
        ("md_simulation", "run_equilibration · run_production", C_EQ),
        ("analyze",   "concat_trajectory · fit_trajectory · analyze_rmsd · analyze_distance · analyze_q_value", C_PROD),
        ("metal",     "detect_metal_ions · parameterize_metal_ion", C_WARN),
        ("literature","pubmed_search · pubmed_fetch", C_EVENT),
        ("slurm",     "submit_job · submit_array_job · check_job · list_tracked_jobs · configure_container", C_AGENT),
        ("node",      "create_node · update_node_status · update_job_params", C_EVENT),
    ]
    row_h = 0.40
    for i, (name, tools, col) in enumerate(groups):
        ry = y0 + i * (row_h + 0.05)
        # pill
        node(slide, MARGIN_X, ry, 1.6, row_h, name, fill=col,
             size=11, bold=True)
        # tools
        text(slide, MARGIN_X + 1.75, ry + 0.06, SLIDE_W - MARGIN_X - 1.75 - 0.3,
             row_h - 0.1, tools, size=11, color=INK)


def figure_antibody_case(slide, y0):
    """SLURM array banner + three parallel mini-DAGs."""
    # dispatcher banner
    banner_y = y0
    filled_rect(slide, MARGIN_X, banner_y, SLIDE_W - 2 * MARGIN_X, 0.55,
                C_WARN)
    text(slide, MARGIN_X, banner_y + 0.12, SLIDE_W - 2 * MARGIN_X, 0.35,
         "submit_array_job(tasks=3)    #SBATCH  --array=0-2",
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, font="Menlo")
    rows = [
        ("batch_1_4m3j_B", "chain B · 119 res", 0),
        ("batch_2_4m3j_A", "chain A · 116 res", 1),
        ("batch_3_4b50_A", "chain A · 122 res", 2),
    ]
    nw, nh = 0.82, 0.42
    xs = [3.05, 3.90, 4.75, 5.60, 6.45, 7.30]
    step_labels_colors = [
        ("source", C_STEP), ("prep", C_STEP), ("solv", C_STEP),
        ("topo", C_STEP), ("eq", C_EQ), ("prod", C_PROD),
    ]
    for job, meta, idx in rows:
        ry = y0 + 0.95 + idx * 0.75
        # left label
        text(slide, MARGIN_X, ry - 0.02, 2.35, 0.25,
             job, size=12, bold=True, color=INK, font="Menlo")
        text(slide, MARGIN_X, ry + 0.22, 2.35, 0.22,
             meta, size=10, color=INK_MUTED, font="Helvetica Neue",
             italic=True)
        for (lbl, col), x in zip(step_labels_colors, xs):
            node(slide, x, ry, nw, nh, lbl, fill=col, size=10)
        for i in range(len(xs) - 1):
            arrow(slide, xs[i] + nw, ry + nh / 2,
                  xs[i + 1], ry + nh / 2, width_pt=1.0)
        # status
        node(slide, 8.30, ry, 1.25, nh, "✓ completed", fill=C_EQ, size=10)
        text(slide, 9.65, ry + 0.06, 2.0, 0.3,
             f"array_task_id = {idx}", size=10, italic=True,
             color=INK_DIM, font="Menlo")
        # vertical connector from banner
        arrow(slide, 6.7, banner_y + 0.55,
              6.7, ry + nh / 2, color=C_WARN, width_pt=1.0)


def figure_deployment(slide, y0):
    """Three columns: plugin / container / HPC."""
    col_w = (SLIDE_W - 2 * MARGIN_X - 0.6) / 3
    cols = [
        ("Plugin", "Claude Code",
         [
             "/plugin install mdclaw@mdclaw",
             "5 skills shipped:",
             "  /md-prepare, /md-equilibration,",
             "  /md-production, /md-analyze, /hpc-run",
             "SessionStart hook auto-pulls SIF",
         ],
         C_STEP),
        ("Container", "ghcr.io/matsunagalab/mdclaw",
         [
             "Docker ~11.4 GB · SIF ~4.6 GB",
             "CUDA 11.8 (driver 520+ floor)",
             "OpenMM 8.2 built from source",
             "AmberTools · Boltz-2 · PyTorch",
             "Works on Pascal → Hopper",
         ],
         C_PROD),
        ("HPC", "slurm_server",
         [
             "submit_job / submit_array_job",
             "check_job · list_tracked_jobs",
             ".mdclaw_jobs.jsonl : job ↔ node",
             "configure_container auto-binds",
             "DAG status synced from sacct",
         ],
         C_AGENT),
    ]
    for i, (tag, sub, items, col) in enumerate(cols):
        cx = MARGIN_X + i * (col_w + 0.3)
        filled_rect(slide, cx, y0, col_w, 4.6, PANEL,
                    line=LINE, line_width_pt=0.75)
        filled_rect(slide, cx, y0, col_w, 0.55, col)
        text(slide, cx, y0 + 0.12, col_w, 0.4, tag,
             size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, cx + 0.25, y0 + 0.75, col_w - 0.5, 0.35,
             sub, size=11, italic=True, color=INK_MUTED,
             font="Menlo")
        tb = _add_textbox(slide, cx + 0.25, y0 + 1.15,
                          col_w - 0.5, 3.3)
        tf = tb.text_frame
        tf.word_wrap = True
        for j, it in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.line_spacing = 1.25
            p.space_after = Pt(5)
            r = p.add_run()
            r.text = it
            r.font.size = Pt(12)
            r.font.color.rgb = INK
            r.font.name = "Menlo" if it.startswith(" ") or \
                it.startswith("/") or "." in it.split()[0] else "Helvetica Neue"


def figure_remd(slide, y0):
    """Coordinator + workers + event log; nested inside an outer-DAG box."""
    # outer DAG band
    filled_rect(slide, MARGIN_X, y0, SLIDE_W - 2 * MARGIN_X, 3.6,
                PANEL, line=INK_DIM, line_width_pt=1.0, dash=True)
    text(slide, MARGIN_X + 0.2, y0 + 0.08, 6.0, 0.3,
         "outer DAG:  remd_001 node  (1 artifact-lineage node)",
         size=11, italic=True, color=INK_MUTED)
    # coordinator
    node(slide, MARGIN_X + 0.3, y0 + 0.75, 2.4, 1.1,
         "Coordinator agent\n(ladder tuning · budget · convergence)",
         fill=C_AGENT, size=12)
    # replicas
    reps = ["rep_000   T = 300 K", "rep_001   T = 310 K", "rep_002   T = 320 K"]
    for i, r in enumerate(reps):
        ry = y0 + 0.55 + i * 0.85
        node(slide, 3.6, ry, 3.0, 0.55, r + "   · worker agent",
             fill=C_STEP, size=11)
        node(slide, 6.85, ry, 2.3, 0.55, "artifacts/trajectory.dcd",
             fill=C_PROD, size=10)
        arrow(slide, 6.6, ry + 0.275, 6.85, ry + 0.275, width_pt=0.9)
        # coordinator to worker
        arrow(slide, MARGIN_X + 0.3 + 2.4, y0 + 1.3,
              3.6, ry + 0.275, color=C_AGENT, width_pt=1.0)
    # event log
    node(slide, 9.45, y0 + 0.85, 3.2, 2.0,
         "events/\nexchange_log.jsonl\n\n(append-only\nexchange events)",
         fill=C_EVENT, size=11)
    for i in range(3):
        ry = y0 + 0.55 + i * 0.85 + 0.275
        arrow(slide, 9.15, ry, 9.45, ry, color=C_EVENT, width_pt=0.9)
    # caption
    text(slide, MARGIN_X, y0 + 3.75, SLIDE_W - 2 * MARGIN_X, 0.3,
         "3-layer separation:  artifact lineage (DAG) · time-ordered communication (events) · private agent state (filesystem subtree)",
         size=11, italic=True, color=INK_MUTED, align=PP_ALIGN.CENTER)


def figure_roadmap(slide, x, y, row_gap=0.58):
    """Level 1-5 staircase."""
    levels = [
        ("L1", "Typed-role edges", "parents: [{node_id, role}]", C_STEP),
        ("L2", "Declarative inputs", "TOPO_SCHEMA = {...}", C_STEP),
        ("L3", "Fan-in aggregators", "umbrella · FEP · REUS", C_EQ),
        ("L4", "Hierarchical compose", "sub-DAG as 1 outer node", C_PROD),
        ("L5", "REMD coupled state", "barrier-sync option (c)", C_AGENT),
    ]
    bh = row_gap - 0.05
    for i, (lv, name, detail, col) in enumerate(levels):
        bx = x
        by = y + i * row_gap
        node(slide, bx, by, 0.55, bh, lv, fill=col, size=12)
        # name in upper half, detail in lower half — non-overlapping
        text(slide, bx + 0.70, by + 0.02, 5.0, 0.26,
             name, size=12, bold=True, color=INK, anchor=MSO_ANCHOR.TOP)
        text(slide, bx + 0.70, by + 0.28, 5.0, 0.26,
             detail, size=9.5, italic=True, color=INK_MUTED,
             font="Menlo", anchor=MSO_ANCHOR.TOP)


# ---------- slides ----------

def slide_title_page(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # left color bar
    filled_rect(s, 0, 0, 0.4, SLIDE_H, ACCENT)
    # project eyebrow
    text(s, 0.9, 1.5, SLIDE_W - 1.3, 0.35,
         "MDCLAW · 進捗報告 · 2026.04",
         size=12, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    # main title
    text(s, 0.9, 2.05, SLIDE_W - 1.3, 1.4,
         "Agent-Orchestrated\nMolecular Dynamics",
         size=44, bold=True, color=INK, line_spacing=1.05)
    # subtitle
    text(s, 0.9, 3.85, SLIDE_W - 1.3, 0.6,
         "A durable DAG harness of artifact lineage, and its applications",
         size=18, italic=True, color=INK_MUTED, line_spacing=1.2)
    # divider
    filled_rect(s, 0.9, 4.75, 3.0, 0.04, ACCENT)
    # author
    text(s, 0.9, 4.95, 10, 0.35,
         "Yasuhiro Matsunaga   ·   Saitama Univ.  &  RIKEN R-CCS (AI4Science Unit)",
         size=15, bold=True, color=INK)
    # date
    text(s, 0.9, 5.35, 8, 0.3,
         "Lab progress meeting  ·  April 20, 2026",
         size=12, color=INK_MUTED)
    # corner stamp (monospace)
    text(s, 0.9, 6.6, SLIDE_W - 1.3, 0.3,
         "mdclaw  v0.x  ·  github.com/matsunagalab/mdclaw  ·  ghcr.io/matsunagalab/mdclaw",
         size=10, color=INK_DIM, font="Menlo")
    return s


def slide_outline(prs):
    s = add_slide(prs, "Outline", "What I will cover today")
    items = [
        ("01", "Motivation", "Why AI agents for MD"),
        ("02", "Primer", "What a DAG is, and why it fits MD workflow state"),
        ("03", "Positioning", "Artifact-lineage DAG, not a LangGraph control flow"),
        ("04", "Implementation", "schema v3 — nodes, events, locks"),
        ("05", "Available CLIs", "11 servers · 51 tools · one mdclaw entrypoint"),
        ("06", "DAG coverage", "Workflows the same harness can express"),
        ("07", "Case study", "Three antibody systems in parallel (batch_1-3)"),
        ("08", "Deployment", "Plugin · Container · HPC"),
        ("09", "Fault tolerance", "A failed node retries from disk — DAG is durable"),
        ("10", "Reproducibility", "What the DAG pins so every run is replayable"),
        ("11", "Ecosystem", "MDDB-AI Hub curator · methods writer · MDEAnalysis · surrogate"),
        ("12", "Correctness (future)", "Evaluator agent · ablation study design"),
        ("13", "Outlook", "Enhanced sampling · educational use · DAG L1-5"),
        ("14", "TODO & asks", "Near-term TODOs and lab requests"),
        ("15", "Try it", "Install, run, give feedback — lab call to action"),
    ]
    row_h = 0.34
    top = BODY_Y + 0.05
    for i, (num, head, sub) in enumerate(items):
        ry = top + i * row_h
        text(s, MARGIN_X + 0.2, ry, 0.9, row_h,
             num, size=12, bold=True, color=ACCENT, font="Menlo",
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, MARGIN_X + 1.2, ry, 3.2, row_h,
             head, size=13, bold=True, color=INK,
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, MARGIN_X + 4.5, ry, 8.1, row_h,
             sub, size=12, color=INK_MUTED,
             anchor=MSO_ANCHOR.MIDDLE)
        filled_rect(s, MARGIN_X + 0.2, ry + row_h - 0.01,
                    SLIDE_W - 2 * MARGIN_X - 0.4, 0.012, LINE)
    return s


def slide_background(prs):
    s = add_slide(prs, "Background",
                  "Motivation — AI for Science, FugakuNEXT, and MD as the testbed")
    bullets(s, MARGIN_X, BODY_Y + 0.05, 7.8, BODY_H, [
        "AI for Science — agents that plan, run, and reason about experiments, not just write code",
        ("Fits the 'auto-research' paradigm (cf. Karpathy): agents as full-stack researchers, not only copilots", 1),
        ("Matsunaga lab × R-CCS AI4Science Unit: prepare for FugakuNEXT-class AI-driven science", 1),
        ("Aligns with two funded lab initiatives:", 1),
        ("NBDC 育成型 MDDB-AI Hub — a Japanese MD data & model commons, federated with European MDDB", 2),
        ("sim2real — transfer-learn nanobody Tm / ddG from MD-generated FEP data into ESM-2", 2),
        "",
        "Why MD is an ideal testbed for agentic auto-research:",
        ("Computationally heavy, already at HPC scale (GPUs today, exascale soon)", 1),
        ("Failure-prone (NaN, OOM, timeout) — exposes durability and recovery edges", 1),
        ("Objective metrics (free energy, RMSD, convergence) — unlike chat benchmarks", 1),
        "",
        "mdclaw plays three roles: (1) data acquisition agent for MDDB-AI Hub, "
        "(2) upstream of sim2real transfer learning, (3) testbed for durable multi-agent auto-research on FugakuNEXT.",
    ], size=12, line_spacing=1.22, space_after=1)
    # right-side: stepping-stone panel
    rx = 8.75
    rw = SLIDE_W - rx - MARGIN_X
    filled_rect(s, rx, BODY_Y + 0.2, rw, 5.1, PANEL,
                line=LINE, line_width_pt=0.5)
    text(s, rx + 0.3, BODY_Y + 0.40, rw - 0.6, 0.3,
         "STEPPING STONE TO FUGAKUNEXT",
         size=11, bold=True, color=ACCENT)
    # three-stage staircase
    stages = [
        ("TODAY",
         "Agents + SLURM GPU cluster\n100 antibody systems completed\nend-to-end, autonomously",
         C_STEP),
        ("NEAR (2026–2028)",
         "MDDB-AI Hub integration (NBDC)\nsim2real: MD → ESM-2 training\nAgent-driven REMD / FEP",
         C_PROD),
        ("FUGAKUNEXT",
         "10⁴–10⁵ durable agent workflows\nrunning concurrently — same DAG,\nsame artifacts",
         C_AGENT),
    ]
    sy = BODY_Y + 0.85
    sh = 1.35
    for i, (tag, body, col) in enumerate(stages):
        y = sy + i * (sh + 0.05)
        filled_rect(s, rx + 0.3, y, rw - 0.6, sh, WHITE,
                    line=col, line_width_pt=1.0)
        filled_rect(s, rx + 0.3, y, 0.12, sh, col)
        text(s, rx + 0.50, y + 0.08, rw - 0.9, 0.3,
             tag, size=11, bold=True, color=col, font="Menlo")
        text(s, rx + 0.50, y + 0.40, rw - 0.9, sh - 0.45,
             body, size=11, color=INK, line_spacing=1.28)
    return s


def slide_dag_primer(prs):
    s = add_slide(prs, "Primer",
                  "What is a DAG, and why it fits MD workflow state")
    bullets(s, MARGIN_X, BODY_Y + 0.05, 12.2, 1.3, [
        "DAG = Directed Acyclic Graph: nodes + directed edges, no cycles",
        "Order and branching are both native — MD pipelines have both",
        "Familiar cousins: git history, Bazel / Make, Snakemake / Nextflow",
        "In mdclaw: a node = one step's artifacts; an edge = 'consumed by'",
    ], size=13, line_spacing=1.25, space_after=2)
    figure_dag_primer(s, y0=4.50)
    return s


def slide_positioning(prs):
    s = add_slide(prs, "Positioning",
                  "mdclaw's DAG is artifact lineage, not LangGraph-style control flow")
    figure_positioning(s, y0=BODY_Y + 0.2)
    # footer takeaway
    text(s, MARGIN_X, 6.55, SLIDE_W - 2 * MARGIN_X, 0.4,
         "Takeaway: durable execution + multi-agent scaling coexist because the DAG is on disk, not in library memory.",
         size=12, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
    return s


def slide_implementation(prs):
    s = add_slide(prs, "Implementation",
                  "DAG harness internals — schema v3: node · event · lock")
    # left: directory tree
    figure_dirtree(s, x=MARGIN_X, y=BODY_Y + 0.1)
    # right: key points
    rx = 7.4
    filled_rect(s, rx, BODY_Y + 0.1, SLIDE_W - rx - MARGIN_X, 4.8,
                PANEL, line=LINE, line_width_pt=0.5)
    text(s, rx + 0.25, BODY_Y + 0.3, 5.0, 0.35,
         "KEY INVARIANTS", size=11, bold=True, color=ACCENT)
    bullets(s, rx + 0.25, BODY_Y + 0.75, SLIDE_W - rx - MARGIN_X - 0.5, 4.0, [
        "one job_dir = one physical system (single-source principle)",
        "node.json holds state · parents · metadata (updated under flock)",
        "artifacts/ = the real scientific data",
        "events/ is append-only — no JSON-array races",
        "progress.json is a thin cached index",
        "",
        "Types: source / prep / solv / topo / eq / prod / analyze",
        "prod branches by continue_from; analyze branches as siblings from one prod",
        "resolve_node_inputs walks ancestors via BFS and injects artifacts",
        "  e.g. an analyze node auto-assembles the prod-lineage DCDs in chronological order",
    ], size=11.5, line_spacing=1.22, space_after=2)
    return s


def slide_clis(prs):
    s = add_slide(prs, "Available CLIs",
                  "11 servers · 51 tools · one mdclaw <tool> entrypoint")
    figure_cli_groups(s, y0=BODY_Y + 0.05)
    text(s, MARGIN_X, 6.75, SLIDE_W - 2 * MARGIN_X, 0.3,
         "Uniform CLI: snake_case → --kebab-case flags · JSON on stdout · --job-dir / --node-id wire any tool into the DAG",
         size=11, italic=True, color=INK_MUTED)
    return s


def slide_coverage(prs):
    s = add_slide(prs, "Coverage",
                  "The same harness expresses a broad set of MD workflows")
    items = [
        ("Linear", "source → prep → solv → topo → eq → prod → analyze", C_STEP),
        ("Branching prod", "prod_001 / prod_002 / prod_003 forked from one eq", C_PROD),
        ("Extension", "continue_from=<prod_id> restores state.xml, keeps timeline", C_PROD),
        ("Mutants", "create_mutated_structure as post-prep variant; downstream solv resolves mutated merged_pdb", C_STEP),
        ("Ligand", "clean_ligand + amber_geostd / GAFF2 fallback", C_EQ),
        ("Membrane", "embed_in_membrane feeds straight into the same DAG", C_EQ),
        ("Large scans", "N job_dirs fanned out via submit_array_job on SLURM", C_WARN),
        ("Analysis subtree",
         "prod → analyze_001 (concat) → sibling analyze_NNN (rmsd / distance / q / fit)",
         C_AGENT),
    ]
    # 2 column grid
    cols = 2
    box_w = (SLIDE_W - 2 * MARGIN_X - 0.4) / cols
    box_h = 0.85
    for idx, (tag, desc, col) in enumerate(items):
        r, c = idx // cols, idx % cols
        bx = MARGIN_X + c * (box_w + 0.4)
        by = BODY_Y + 0.1 + r * (box_h + 0.2)
        filled_rect(s, bx, by, box_w, box_h, PANEL,
                    line=LINE, line_width_pt=0.5)
        filled_rect(s, bx, by, 0.15, box_h, col)
        text(s, bx + 0.3, by + 0.1, box_w - 0.45, 0.35,
             tag, size=14, bold=True, color=INK)
        text(s, bx + 0.3, by + 0.42, box_w - 0.45, 0.4,
             desc, size=12, color=INK_MUTED, font="Menlo")
    text(s, MARGIN_X, 6.75, SLIDE_W - 2 * MARGIN_X, 0.3,
         "Next wave (slide 10): fan-in aggregators for FEP / REMD inside the same DAG model",
         size=11, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
    return s


def slide_case_study(prs):
    s = add_slide(prs, "Case study",
                  "Antibody screening — 3 systems verified locally, 100 systems completed on HPC")
    # lead sentence
    text(s, MARGIN_X, BODY_Y + 0.0, SLIDE_W - 2 * MARGIN_X, 0.35,
         "Targets: PDB 4M3J (chains A/B) and 4B50 (chain A), and 100 antibody variants on the GPU cluster.  "
         "Shared config: ff19SB · OPC · explicit · autonomous.",
         size=12, color=INK_MUTED)
    figure_antibody_case(s, y0=BODY_Y + 0.5)
    # scale callout (big number panel on the right)
    cx = 11.2
    cy = 3.15
    filled_rect(s, cx, cy, 1.85, 2.15, ACCENT_LIGHT,
                line=ACCENT, line_width_pt=1.0)
    text(s, cx, cy + 0.2, 1.85, 0.4, "AT SCALE",
         size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    text(s, cx, cy + 0.55, 1.85, 1.0, "100",
         size=56, bold=True, color=ACCENT_DARK, align=PP_ALIGN.CENTER,
         line_spacing=1.0)
    text(s, cx, cy + 1.45, 1.85, 0.35, "antibody systems",
         size=11, bold=True, color=INK, align=PP_ALIGN.CENTER)
    text(s, cx, cy + 1.75, 1.85, 0.35,
         "completed on\nSLURM + GPU cluster",
         size=10, italic=True, color=INK_MUTED, align=PP_ALIGN.CENTER,
         line_spacing=1.2)
    # status footer
    text(s, MARGIN_X, 6.6, SLIDE_W - 2 * MARGIN_X, 0.4,
         "End-to-end with no human intervention per job:  each node.json carries slurm_job_id and array_task_id, "
         "failures auto-tracked in .mdclaw_jobs.jsonl.",
         size=12, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
    return s


def slide_deployment(prs):
    s = add_slide(prs, "Deployment",
                  "Plugin + container + HPC — usable anywhere, by anyone")
    figure_deployment(s, y0=BODY_Y + 0.1)
    text(s, MARGIN_X, 6.75, SLIDE_W - 2 * MARGIN_X, 0.3,
         "Identical behaviour from Claude Code CLI · Anthropic API · Cursor · Windsurf.",
         size=11, italic=True, color=INK_MUTED, align=PP_ALIGN.CENTER)
    return s


def figure_node_json(slide, x, y, w, h):
    """Stylized node.json preview with color-coded lines mapping to mechanisms."""
    filled_rect(slide, x, y, w, h, PANEL, line=LINE, line_width_pt=0.5)
    # header strip
    filled_rect(slide, x, y, w, 0.4, INK)
    text(slide, x + 0.2, y + 0.08, w - 0.4, 0.3,
         "nodes/prod_003/node.json", size=11, bold=True, color=WHITE,
         font="Menlo")
    lines = [
        ('{',                                              INK, False),
        ('  "id": "prod_003",',                            INK, False),
        ('  "type": "prod",',                              INK, False),
        ('  "parents": ["eq_001"],',                       C_STEP, True),
        ('  "status": "completed",',                       INK, False),
        ('  "metadata": {',                                INK, False),
        ('    "random_seed": 42,',                         C_AGENT, True),
        ('    "hmr": true,',                               C_AGENT, True),
        ('    "timestep_fs": 4.0,',                        C_AGENT, True),
        ('    "simulation_time_ns": 100,',                 INK_MUTED, False),
        ('    "start_step": 12500000,',                    INK_MUTED, False),
        ('    "final_step": 37500000,',                    INK_MUTED, False),
        ('    "forcefield": "ff19SB",',                    C_EQ, True),
        ('    "water_model": "opc",',                      C_EQ, True),
        ('    "tool_version": "mdclaw 0.9.3",',            C_PROD, True),
        ('    "container_digest": "sha256:9f3a…",',        C_PROD, True),
        ('    "input_sha256": {',                          C_WARN, True),
        ('      "system_xml": "3b7c…",',                   C_WARN, True),
        ('      "topology_pdb": "a14e…",',                 C_WARN, True),
        ('      "state_xml": "9d22…"',                     C_WARN, True),
        ('    }',                                          INK, False),
        ('  }',                                            INK, False),
        ('}',                                              INK, False),
    ]
    tb = _add_textbox(slide, x + 0.2, y + 0.5, w - 0.4, h - 0.6)
    tf = tb.text_frame
    tf.word_wrap = False
    for i, (line, col, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.12
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = line
        r.font.name = "Menlo"
        r.font.size = Pt(10)
        r.font.color.rgb = col
        r.font.bold = bold


def slide_reproducibility(prs):
    s = add_slide(prs, "Reproducibility",
                  "Reproducibility by construction — every node pins its inputs")
    # left: stylized node.json
    figure_node_json(s, x=MARGIN_X, y=BODY_Y + 0.1, w=6.2, h=5.1)

    # right: mechanisms as colored cards (2x3 grid)
    rx = MARGIN_X + 6.2 + 0.3
    rw = SLIDE_W - rx - MARGIN_X
    mech = [
        ("DAG lineage",
         "parents: [...] chain fixes every input",
         "Re-runs walk the same ancestors",
         C_STEP),
        ("Deterministic run config",
         "random_seed · HMR · timestep_fs",
         "Written before run, not after",
         C_AGENT),
        ("Canonicalized parameters",
         "ff / water / ions normalized",
         "No 'opc' vs 'OPC' mismatch",
         C_EQ),
        ("Code & binary pinned",
         "tool_version + container_digest",
         "SIF sha256 tied to plugin version",
         C_PROD),
        ("Content-addressed inputs",
         "SHA256 of system.xml / topology.pdb / state.xml",
         "Detects silent upstream edits",
         C_WARN),
        ("Append-only events + flock",
         "events/ never rewritten, atomic status",
         "Audit trail survives crashes",
         C_EVENT),
    ]
    cols = 2
    rows = 3
    pad = 0.12
    cw = (rw - pad * (cols - 1)) / cols
    ch = (5.1 - pad * (rows - 1)) / rows
    for idx, (head, code, why, col) in enumerate(mech):
        r, c = idx // cols, idx % cols
        cx = rx + c * (cw + pad)
        cy = BODY_Y + 0.1 + r * (ch + pad)
        filled_rect(s, cx, cy, cw, ch, PANEL, line=LINE, line_width_pt=0.5)
        filled_rect(s, cx, cy, 0.12, ch, col)
        text(s, cx + 0.22, cy + 0.12, cw - 0.35, 0.3,
             head, size=12, bold=True, color=INK)
        text(s, cx + 0.22, cy + 0.45, cw - 0.35, 0.3,
             code, size=10, color=col, font="Menlo", bold=True)
        text(s, cx + 0.22, cy + 0.78, cw - 0.35, 0.55,
             why, size=10, italic=True, color=INK_MUTED,
             line_spacing=1.25)

    # footer takeaway
    text(s, MARGIN_X, 6.65, SLIDE_W - 2 * MARGIN_X, 0.35,
         "Rerun mdclaw on the same job_dir, get the same trajectory — "
         "bit-identical on same GPU, statistically identical across GPUs.",
         size=12, italic=True, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    return s


def figure_retry(slide, x, y):
    """DAG showing a failed attempt and a successful retry on the same node path."""
    nw, nh = 1.1, 0.5
    # baseline chain: eq → prod (two attempts)
    node(slide, x, y + 0.8, nw, nh, "eq_001", fill=C_EQ)
    # attempt 1: failed
    attempt1_x = x + nw + 0.8
    attempt1_y = y
    node(slide, attempt1_x, attempt1_y, nw + 0.4, nh + 0.1,
         "prod_001\nattempt 1", fill=INK_DIM, size=10)
    # cross-out / failure tag
    text(slide, attempt1_x - 0.35, attempt1_y + 0.75, nw + 1.1, 0.3,
         "✗ FAILED",
         size=10, bold=True, color=C_AGENT, align=PP_ALIGN.CENTER)
    text(slide, attempt1_x - 0.35, attempt1_y + 1.02, nw + 1.1, 0.3,
         "NaN at step 1.2e7",
         size=9, italic=True, color=INK_MUTED, align=PP_ALIGN.CENTER,
         font="Menlo")
    # attempt 2: success
    attempt2_x = x + nw + 0.8
    attempt2_y = y + 1.7
    node(slide, attempt2_x, attempt2_y, nw + 0.4, nh + 0.1,
         "prod_001\nattempt 2", fill=C_PROD, size=10)
    text(slide, attempt2_x - 0.35, attempt2_y + 0.72, nw + 1.1, 0.3,
         "✓ COMPLETED",
         size=10, bold=True, color=C_EQ, align=PP_ALIGN.CENTER)
    text(slide, attempt2_x - 0.35, attempt2_y + 0.98, nw + 1.1, 0.3,
         "restarted from eq state.xml, new seed",
         size=9, italic=True, color=INK_MUTED, align=PP_ALIGN.CENTER,
         font="Menlo")
    # arrows
    arrow(slide, x + nw, y + 0.8 + nh / 2,
          attempt1_x, attempt1_y + nh / 2 + 0.05, width_pt=1.4)
    arrow(slide, x + nw, y + 0.8 + nh / 2,
          attempt2_x, attempt2_y + nh / 2 + 0.05, width_pt=1.4)
    # retry arrow between attempts
    arrow(slide, attempt1_x + (nw + 0.4) / 2, attempt1_y + nh + 0.15,
          attempt2_x + (nw + 0.4) / 2, attempt2_y - 0.05,
          color=C_AGENT, width_pt=1.6)
    text(slide, attempt1_x + nw - 0.2, attempt1_y + 1.35, 2.0, 0.3,
         "retry", size=11, bold=True, italic=True, color=C_AGENT)
    # status transitions below
    ty = y + 2.85
    statuses = [("queued", INK_DIM), ("running", C_STEP),
                ("failed", C_AGENT), ("queued", INK_DIM),
                ("running", C_STEP), ("completed", C_EQ)]
    tx = x + 0.1
    tw = 0.85
    for i, (st, col) in enumerate(statuses):
        node(slide, tx + i * (tw + 0.1), ty, tw, 0.32, st,
             fill=col, size=9)
        if i < len(statuses) - 1:
            arrow(slide, tx + i * (tw + 0.1) + tw, ty + 0.16,
                  tx + (i + 1) * (tw + 0.1), ty + 0.16, width_pt=0.9)
    text(slide, x, ty + 0.42, 7.0, 0.25,
         "node.json status transitions, atomically updated under flock",
         size=10, italic=True, color=INK_MUTED, font="Menlo")


def slide_fault_tolerance(prs):
    s = add_slide(prs, "Fault tolerance",
                  "The DAG is durable — a failed node restarts from disk")
    # left: retry figure
    figure_retry(s, x=MARGIN_X + 0.2, y=BODY_Y + 0.2)
    # right: explanation
    rx = 7.2
    filled_rect(s, rx, BODY_Y + 0.1, SLIDE_W - rx - MARGIN_X, 5.0,
                PANEL, line=LINE, line_width_pt=0.5)
    text(s, rx + 0.25, BODY_Y + 0.3, 5.5, 0.3,
         "WHAT WE RECOVER FROM", size=11, bold=True, color=ACCENT)
    bullets(s, rx + 0.25, BODY_Y + 0.7, SLIDE_W - rx - MARGIN_X - 0.5, 1.8, [
        "OpenMM NaN / particle clash",
        "GPU OOM, CPU OOM",
        "SLURM timeout / node eviction",
        "Transient filesystem or network errors",
        "User-triggered restart with new params",
    ], size=11, line_spacing=1.2, space_after=2, indent_color=ACCENT)

    text(s, rx + 0.25, BODY_Y + 2.85, 5.5, 0.3,
         "RECOVERY MECHANICS", size=11, bold=True, color=ACCENT)
    bullets(s, rx + 0.25, BODY_Y + 3.25, SLIDE_W - rx - MARGIN_X - 0.5, 1.8, [
        "update_node_status writes failed + stderr_tail (under flock)",
        "Agent re-queues the node, or branches via continue_from",
        "resolve_node_inputs re-loads state.xml from nearest ancestor",
        "events/ keeps the full attempt history — audit trail preserved",
    ], size=11, line_spacing=1.2, space_after=2, indent_color=ACCENT)
    # bottom takeaway
    text(s, MARGIN_X, 6.65, SLIDE_W - 2 * MARGIN_X, 0.35,
         "All state lives on disk.  Agent RAM is disposable — any process can pick up the DAG and continue.",
         size=12, italic=True, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)
    return s


def figure_ablation(slide, x, y, w, h):
    """2x2 matrix: DAG × Evaluator."""
    # outer border
    filled_rect(slide, x, y, w, h, WHITE, line=LINE, line_width_pt=0.75)
    # column headers
    col_w = (w - 1.6) / 2
    filled_rect(slide, x + 1.6, y, col_w, 0.5, INK_DIM)
    filled_rect(slide, x + 1.6 + col_w, y, col_w, 0.5, ACCENT)
    text(slide, x + 1.6, y + 0.1, col_w, 0.3, "No DAG",
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(slide, x + 1.6 + col_w, y + 0.1, col_w, 0.3, "With DAG (mdclaw)",
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # row headers
    row_h = (h - 0.5) / 2
    filled_rect(slide, x, y + 0.5, 1.6, row_h, INK_DIM)
    filled_rect(slide, x, y + 0.5 + row_h, 1.6, row_h, ACCENT)
    text(slide, x, y + 0.5, 1.6, row_h, "No\nevaluator",
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    text(slide, x, y + 0.5 + row_h, 1.6, row_h, "With\nevaluator",
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    # cells
    cells = [
        (0, 0, "classical scripts", "baseline", INK_DIM),
        (0, 1, "mdclaw today", "today", C_STEP),
        (1, 0, "scripts + auto-QA", "checks only", INK_DIM),
        (1, 1, "mdclaw + evaluator", "proposed", C_AGENT),
    ]
    for r, c, head, sub, col in cells:
        cx = x + 1.6 + c * col_w
        cy = y + 0.5 + r * row_h
        filled_rect(slide, cx + 0.03, cy + 0.03, col_w - 0.06, row_h - 0.06,
                    PANEL)
        text(slide, cx + 0.1, cy + 0.12, col_w - 0.2, 0.35,
             head, size=13, bold=True, color=col)
        text(slide, cx + 0.1, cy + 0.48, col_w - 0.2, 0.3,
             sub, size=10, italic=True, color=INK_MUTED)


def slide_evaluator_ablation(prs):
    s = add_slide(prs, "Correctness (future)",
                  "Evaluator agent & ablation study — how we plan to measure the system")
    # left: evaluator agent panel
    lx = MARGIN_X
    lw = 6.0
    filled_rect(s, lx, BODY_Y + 0.1, lw, 5.0, PANEL,
                line=LINE, line_width_pt=0.5)
    filled_rect(s, lx, BODY_Y + 0.1, lw, 0.65, C_AGENT)
    text(s, lx + 0.2, BODY_Y + 0.15, lw - 0.4, 0.32,
         "MD evaluator = MDDB-AI Hub quality gate",
         size=13, bold=True, color=WHITE)
    text(s, lx + 0.2, BODY_Y + 0.45, lw - 0.4, 0.25,
         "same agent judges correctness AND ingestion-worthiness",
         size=10, italic=True, color=WHITE)
    text(s, lx + 0.3, BODY_Y + 0.90, lw - 0.6, 0.3,
         "Two axes of evaluation:",
         size=11, color=INK, italic=True)
    # Axis 1 — correctness
    text(s, lx + 0.3, BODY_Y + 1.25, lw - 0.6, 0.25,
         "① CORRECTNESS  — did the MD actually run right?",
         size=11, bold=True, color=C_AGENT)
    bullets(s, lx + 0.45, BODY_Y + 1.52, lw - 0.75, 1.2, [
        "Numerical stability — no NaN, bounded energy drift",
        "Thermostat / barostat — T, P, density stable",
        "Conservation — momentum ~ 0, volume bounded",
        "Protocol compliance — seed / HMR / step count honored",
    ], size=10.5, line_spacing=1.2, space_after=1, indent_color=C_AGENT)
    # Axis 2 — data quality
    text(s, lx + 0.3, BODY_Y + 2.80, lw - 0.6, 0.25,
         "② DATA QUALITY — is the distribution worth keeping?",
         size=11, bold=True, color=C_STEP)
    bullets(s, lx + 0.45, BODY_Y + 3.07, lw - 0.75, 1.2, [
        "Sampling coverage — conformational entropy, block-avg",
        "Convergence — autocorrelation, KS across halves",
        "Structural sanity — RMSD / RMSF in expected range",
        "MDDB-AI Hub ingestion criteria (min length, …)",
    ], size=10.5, line_spacing=1.2, space_after=1, indent_color=C_STEP)
    # output contract
    verdicts = [("PASS", C_EQ), ("WARN", C_WARN), ("FAIL", C_AGENT)]
    vy = BODY_Y + 4.45
    text(s, lx + 0.3, vy, 1.5, 0.35,
         "Output  →", size=11, bold=True, color=INK,
         anchor=MSO_ANCHOR.MIDDLE)
    vx = lx + 1.8
    for i, (v, c) in enumerate(verdicts):
        node(s, vx + i * 1.15, vy, 1.0, 0.35, v, fill=c, size=11)
    text(s, lx + 0.3, vy + 0.42, lw - 0.6, 0.25,
         "FAIL → retry · WARN → log · PASS → mark node ingestable",
         size=10, italic=True, color=INK_MUTED, font="Menlo")

    # right: ablation matrices (core + meta)
    rx = lx + lw + 0.3
    rw = SLIDE_W - rx - MARGIN_X
    text(s, rx, BODY_Y + 0.15, rw, 0.3,
         "CORE ABLATION — DAG × EVALUATOR",
         size=11, bold=True, color=ACCENT)
    figure_ablation(s, x=rx, y=BODY_Y + 0.5, w=rw, h=2.05)

    # meta-ablation: LLM × harness sweep
    ma_y = BODY_Y + 2.75
    text(s, rx, ma_y, rw, 0.3,
         "META-ABLATION — LLM × HARNESS SWEEP",
         size=11, bold=True, color=ACCENT)
    # two columns: LLMs and Harnesses
    sub_w = (rw - 0.3) / 2
    filled_rect(s, rx, ma_y + 0.35, sub_w, 1.8, PANEL,
                line=LINE, line_width_pt=0.5)
    filled_rect(s, rx + sub_w + 0.3, ma_y + 0.35, sub_w, 1.8, PANEL,
                line=LINE, line_width_pt=0.5)
    text(s, rx + 0.15, ma_y + 0.45, sub_w - 0.3, 0.3,
         "LLM", size=10, bold=True, color=INK_MUTED)
    bullets(s, rx + 0.15, ma_y + 0.75, sub_w - 0.3, 1.4, [
        "Claude Opus 4.7",
        "Claude Sonnet 4.6",
        "Claude Haiku 4.5",
        "GPT-5 / o-series",
        "Gemini 3 Pro",
    ], size=10, line_spacing=1.15, space_after=1,
       indent_color=C_STEP, font="Menlo")
    text(s, rx + sub_w + 0.45, ma_y + 0.45, sub_w - 0.3, 0.3,
         "Agent harness", size=10, bold=True, color=INK_MUTED)
    bullets(s, rx + sub_w + 0.45, ma_y + 0.75, sub_w - 0.3, 1.4, [
        "Claude Code",
        "Cursor",
        "OpenCode",
        "Codex CLI",
        "pi / Amp / …",
    ], size=10, line_spacing=1.15, space_after=1,
       indent_color=C_AGENT, font="Menlo")

    # metrics line
    text(s, rx, BODY_Y + 4.7, rw, 0.35,
         "Metrics: wall-time · recovery rate · param errors caught · "
         "human interventions · reproducibility",
         size=10, italic=True, color=INK_MUTED)
    return s


def slide_ecosystem(prs):
    s = add_slide(prs, "Ecosystem",
                  "Agentic ecosystem — metadata, authoring, analysis, surrogate")
    panels = [
        ("MDDB-AI Hub curator",
         "Register every DAG  (NBDC)",
         "DAG state ≈ ideal\n"
         "MD metadata:\n"
         "  ff, water, ions, box,\n"
         "  HMR, timestep, seed,\n"
         "  ancestor lineage.\n\n"
         "Curator reads node.json\n"
         "chain → POSTs to the\n"
         "MDDB-AI Hub (federated\n"
         "with European MDDB /\n"
         "IRB Barcelona).\n\n"
         "On HPCI共用 1.5 PB.",
         C_STEP),
        ("Paper-methods writer",
         "Draft the Methods",
         "Input:  progress.json +\n"
         "node.json chain + events/\n"
         "+ artifacts/ headers.\n\n"
         "Output: Markdown / LaTeX\n"
         "Methods section with every\n"
         "parameter pinned —\n"
         "reproducibility by\n"
         "construction.\n\n"
         "Downstream: citation,\n"
         "consistency checks vs\n"
         "prior lab papers.",
         C_PROD),
        ("MDEAnalysis CLI",
         "Explain one run",
         "mdclaw analyze explain\n"
         "    <prod_id>\n\n"
         "Walks prod → eq →\n"
         "topo → solv → …\n"
         "Stitches segments into\n"
         "one timeline, annotates\n"
         "seeds / restarts / HMR.\n\n"
         "Agent output: 'this 100 ns\n"
         "was restarted from\n"
         "prod_002 with seed 42\n"
         "after a NaN at 1.2e7.'",
         C_EQ),
        ("MD surrogate CLI",
         "Swap MD for ML surrogate",
         "mdclaw surrogate run\n"
         "    <job_dir>\n\n"
         "Replaces run_production\n"
         "with a neural surrogate\n"
         "(Boltz-2, AlphaFlow,\n"
         "MACE-MD, …).\n\n"
         "Same DAG, same node\n"
         "types, new artifact kind.\n\n"
         "Use case: screen 10³ var-\n"
         "iants in minutes, then\n"
         "promote finalists to full MD.",
         C_AGENT),
    ]
    n = len(panels)
    col_w = (SLIDE_W - 2 * MARGIN_X - 0.2 * (n - 1)) / n
    for i, (head, sub, body, col) in enumerate(panels):
        cx = MARGIN_X + i * (col_w + 0.2)
        y = BODY_Y + 0.1
        filled_rect(s, cx, y, col_w, 4.95, PANEL,
                    line=LINE, line_width_pt=0.5)
        filled_rect(s, cx, y, col_w, 0.7, col)
        text(s, cx + 0.15, y + 0.08, col_w - 0.3, 0.32,
             head, size=13, bold=True, color=WHITE)
        text(s, cx + 0.15, y + 0.40, col_w - 0.3, 0.28,
             sub, size=9.5, italic=True, color=WHITE)
        text(s, cx + 0.2, y + 0.9, col_w - 0.4, 3.95,
             body, size=10.5, color=INK, line_spacing=1.30,
             font="Helvetica Neue")
    # footer
    text(s, MARGIN_X, 6.75, SLIDE_W - 2 * MARGIN_X, 0.3,
         "The DAG is not only execution state — it is the unit we register, cite, analyze, and swap engines on.",
         size=11, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
    return s


def slide_aim(prs):
    s = add_slide(prs, "Outlook",
                  "Enhanced sampling · educational use · DAG relaxation L1-5")
    # left column: REMD figure
    fig_x = MARGIN_X
    fig_y = BODY_Y + 0.1
    fig_w = 7.60
    fig_h = 2.95
    # outer DAG band
    filled_rect(s, fig_x, fig_y, fig_w, fig_h,
                PANEL, line=INK_DIM, line_width_pt=1.0, dash=True)
    text(s, fig_x + 0.15, fig_y + 0.06, 7.0, 0.28,
         "enhanced sampling: Sub-Agent REMD inside a single outer DAG node",
         size=10, italic=True, color=INK_MUTED)
    # coordinator
    node(s, fig_x + 0.25, fig_y + 0.55, 1.9, 1.1,
         "Coordinator agent\n(ladder · budget)",
         fill=C_AGENT, size=11)
    reps = ["rep_000  T=300 K", "rep_001  T=310 K", "rep_002  T=320 K"]
    for i, r in enumerate(reps):
        ry = fig_y + 0.45 + i * 0.70
        node(s, fig_x + 2.40, ry, 2.1, 0.45, r + " · worker",
             fill=C_STEP, size=10)
        node(s, fig_x + 4.65, ry, 1.55, 0.45, "trajectory.dcd",
             fill=C_PROD, size=10)
        arrow(s, fig_x + 4.50, ry + 0.225, fig_x + 4.65, ry + 0.225,
              width_pt=0.9)
        arrow(s, fig_x + 2.15, fig_y + 1.10, fig_x + 2.40, ry + 0.225,
              color=C_AGENT, width_pt=1.0)
    # event log (wider so it can breathe)
    node(s, fig_x + 6.35, fig_y + 0.55, 1.15, 2.05,
         "events/\nexchange_log\n.jsonl", fill=C_EVENT, size=10)
    for i in range(3):
        ry = fig_y + 0.45 + i * 0.70 + 0.225
        arrow(s, fig_x + 6.20, ry, fig_x + 6.35, ry,
              color=C_EVENT, width_pt=0.9)

    # right column: relaxation roadmap
    rx = fig_x + fig_w + 0.25
    rw = SLIDE_W - rx - MARGIN_X
    text(s, rx, BODY_Y + 0.1, rw, 0.3,
         "DAG RELAXATION ROADMAP", size=11, bold=True, color=ACCENT)
    figure_roadmap(s, x=rx, y=BODY_Y + 0.45, row_gap=0.52)

    # bottom band: 3 future horizons (below both figure and roadmap)
    by = BODY_Y + 0.1 + max(fig_h, 0.45 + 5 * 0.52) + 0.3
    horizons = [
        ("Enhanced sampling",
         "Umbrella · FEP · weighted ensemble fit the DAG natively (branch + fan-in).  "
         "REMD needs extra work — coupled state, L5 barrier-sync.",
         C_AGENT),
        ("Educational use",
         "Step-by-step natural-language tutorials generated from real DAG state — students see every decision",
         C_EQ),
        ("Sim2real ML pipeline",
         "MDDB-AI Hub DAGs feed sim2real training — FEP/ddG traj. → ESM-2 fine-tune for Tm / binding prediction",
         C_PROD),
    ]
    col_w = (SLIDE_W - 2 * MARGIN_X - 0.4) / 3
    for i, (h, body, col) in enumerate(horizons):
        cx = MARGIN_X + i * (col_w + 0.2)
        filled_rect(s, cx, by, col_w, 1.1, WHITE,
                    line=col, line_width_pt=1.0)
        filled_rect(s, cx, by, 0.15, 1.1, col)
        text(s, cx + 0.3, by + 0.1, col_w - 0.45, 0.3,
             h, size=12, bold=True, color=col)
        text(s, cx + 0.3, by + 0.42, col_w - 0.45, 0.7,
             body, size=10.5, color=INK_MUTED, line_spacing=1.25)
    return s


def slide_call_to_action(prs):
    s = add_slide(prs, "Try it & feedback",
                  "Please try mdclaw — your feedback shapes the next iteration")
    # two big zones: "Try it" / "Feedback wanted" / "Where to report"
    # Left: quick start
    lx = MARGIN_X
    lw = 6.3
    ly = BODY_Y + 0.1
    filled_rect(s, lx, ly, lw, 5.0, PANEL, line=LINE, line_width_pt=0.5)
    filled_rect(s, lx, ly, lw, 0.55, ACCENT)
    text(s, lx + 0.25, ly + 0.13, lw - 0.5, 0.35,
         "TRY IT IN 3 LINES", size=13, bold=True, color=WHITE)
    # code block
    code_lines = [
        "# install the plugin (Claude Code)",
        "/plugin marketplace add matsunagalab/mdclaw",
        "/plugin install mdclaw@mdclaw",
        "",
        "# try it on any PDB ID",
        "/mdclaw:md-prepare       1AKE chain A, defaults",
        "/mdclaw:md-equilibration job_a1b2c3d4",
        "/mdclaw:md-production    job_a1b2c3d4, 100 ns",
        "/mdclaw:md-analyze       job_a1b2c3d4",
        "",
        "# or on the HPC — one line to fan out N variants",
        "/mdclaw:hpc-run submit 100 ns of 1AKE on gpu01",
    ]
    tb = _add_textbox(s, lx + 0.3, ly + 0.75, lw - 0.6, 3.8)
    tf = tb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Menlo"
        r.font.size = Pt(11)
        if line.startswith("#"):
            r.font.color.rgb = INK_DIM
            r.font.italic = True
        elif line.startswith("/"):
            r.font.color.rgb = ACCENT_DARK
            r.font.bold = True
        else:
            r.font.color.rgb = INK
    # last note
    text(s, lx + 0.3, ly + 4.55, lw - 0.6, 0.35,
         "Works the same from Claude Code, API, Cursor, Windsurf.",
         size=11, italic=True, color=INK_MUTED, align=PP_ALIGN.LEFT)

    # Right: feedback wanted + where to report
    rx = lx + lw + 0.3
    rw = SLIDE_W - rx - MARGIN_X

    # top card: what feedback
    filled_rect(s, rx, ly, rw, 3.1, PANEL, line=LINE, line_width_pt=0.5)
    filled_rect(s, rx, ly, rw, 0.55, C_PROD)
    text(s, rx + 0.25, ly + 0.13, rw - 0.5, 0.35,
         "FEEDBACK I ESPECIALLY WANT", size=13, bold=True, color=WHITE)
    bullets(s, rx + 0.3, ly + 0.75, rw - 0.6, 2.3, [
        "Parameter defaults that don't match your science",
        "Tools missing for your workflow — what would you add?",
        "Error messages that were confusing or unhelpful",
        "DAG node types you wish existed (fep, umbrella, …)",
        "Systems where the agent made the wrong judgement call",
    ], size=11, line_spacing=1.3, space_after=3, indent_color=C_PROD)

    # bottom card: where to report
    filled_rect(s, rx, ly + 3.3, rw, 1.7, PANEL, line=LINE, line_width_pt=0.5)
    filled_rect(s, rx, ly + 3.3, rw, 0.55, C_AGENT)
    text(s, rx + 0.25, ly + 3.43, rw - 0.5, 0.35,
         "WHERE TO REPORT", size=13, bold=True, color=WHITE)
    bullets(s, rx + 0.3, ly + 4.05, rw - 0.6, 1.0, [
        "GitHub issues — github.com/matsunagalab/mdclaw",
        "Lab Slack #mdclaw — daily thread",
        "Or just grab me in the corridor",
    ], size=11, line_spacing=1.3, space_after=2, indent_color=C_AGENT)

    # footer
    text(s, MARGIN_X, 6.65, SLIDE_W - 2 * MARGIN_X, 0.35,
         "The tool matures in proportion to the variety of systems we throw at it — please throw yours.",
         size=12, italic=True, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    return s


def slide_todo(prs):
    s = add_slide(prs, "TODO & asks", "Near-term work and asks to the lab")
    # three panels: TODO · papers · asks
    panels = [
        ("NEAR-TERM TODO", C_EQ, [
            "Run /md-analyze on the 100 antibody DAGs (analyze node type now live — concat + rmsd / distance / q)",
            "Add secondary-structure / H-bond / PCA tools to the [analyze] server",
            "Level-1 typed-role edges in _node.py (backward-compatible)",
            "MDDB-AI Hub integration — register finished jobs to the NBDC database",
            "Sim2real data pipeline — dump FEP/ddG artifacts to sim2real training format",
        ]),
        ("IDEAS TO EXPLORE (memo)", C_WARN, [
            "Alchemical topology tool — wrap parmed / openmmtools.alchemy to build λ-interpolated system.xml  →  unlocks FEP natively on the DAG",
            "MD surrogate CLI — swap run_production for Boltz-2 / AlphaFlow / MACE-MD for fast screening",
            "MDEAnalysis CLI — walk DAG from leaf upward, narrate one trajectory's provenance",
        ]),
        ("PAPER STRATEGY (TENTATIVE)", C_PROD, [
            "Paper 1: L1-3 impl + MD case study (antibody / mutant scan)",
            "Paper 2: L5(c) + Sub-Agent REMD — theory, experiments, logs",
        ]),
        ("ASKS TO THE LAB", C_AGENT, [
            "GPU budget — REMD needs N replicas × several days, queue policy?",
            "Collaboration seeds — antibody / mutant screening welcome",
        ]),
    ]
    top = BODY_Y + 0.05
    panel_h = 1.35
    gap = 0.08
    for i, (title, col, items) in enumerate(panels):
        y = top + i * (panel_h + gap)
        filled_rect(s, MARGIN_X, y, SLIDE_W - 2 * MARGIN_X, panel_h, PANEL,
                    line=LINE, line_width_pt=0.5)
        filled_rect(s, MARGIN_X, y, 0.15, panel_h, col)
        text(s, MARGIN_X + 0.3, y + 0.08, 6.0, 0.3,
             title, size=11, bold=True, color=col)
        bullets(s, MARGIN_X + 0.3, y + 0.38,
                SLIDE_W - 2 * MARGIN_X - 0.5, panel_h - 0.42,
                items, size=11, line_spacing=1.2, space_after=1,
                indent_color=col)
    return s


# ---------- main ----------

def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # build
    slide_title_page(prs)
    slide_outline(prs)
    slide_background(prs)
    slide_dag_primer(prs)
    slide_positioning(prs)
    slide_implementation(prs)
    slide_clis(prs)
    slide_coverage(prs)
    slide_case_study(prs)
    slide_deployment(prs)
    slide_fault_tolerance(prs)
    slide_reproducibility(prs)
    slide_ecosystem(prs)
    slide_evaluator_ablation(prs)
    slide_aim(prs)
    slide_todo(prs)
    slide_call_to_action(prs)

    # apply chrome (context bar, page number, title) to non-title slides
    total = len(prs.slides)
    for i, s in enumerate(prs.slides):
        if i == 0:
            continue
        ctx = getattr(s, "_mdclaw_context", "")
        title = getattr(s, "_mdclaw_title", "")
        subtitle = getattr(s, "_mdclaw_subtitle", None)
        chrome(s, ctx, i + 1, total)
        slide_title(s, title, subtitle)

    prs.save(str(OUTPUT))
    print(f"saved: {OUTPUT}  ({OUTPUT.stat().st_size} bytes, "
          f"{len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
